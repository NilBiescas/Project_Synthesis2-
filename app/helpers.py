from doctr.io import DocumentFile
from doctr.models import ocr_predictor
from pptx.util import Pt
from pptx.enum.text import PP_ALIGN
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_AUTO_SIZE
from pptx.enum.text import MSO_ANCHOR
import matplotlib.pyplot as plt
import aspose.pydrawing as draw
import aspose.slides as slides
import os


def docTR_inference(png_path, device="cpu"):
    model = ocr_predictor(pretrained=True)
    model.to(device)
    doc = DocumentFile.from_images(png_path)
    result = model(doc)
    return result, doc

    
def add_line_in_bounding_boxes(presentation, slide_index, text, coordinates):
    slide = presentation.slides[slide_index]
    left, top, width, height = coordinates
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    # tf.word_wrap = True
    # print(0.75*coordinates[3])
    tf.paragraphs[0].font.size = Pt(int(0.00007*coordinates[3]))
    tf.paragraphs[0].text = text
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT

def plot_slide(presentation, slide, out_dir):

    with slides.Presentation(presentation) as pres:
        # Access the first slide
        sld = pres.slides[slide]

        # User defined dimension
        desiredX = 210*4
        desiredY = 297*4

        # Getting scaled value  of X and Y
        ScaleX = (1.0 / pres.slide_size.size.width) * desiredX
        ScaleY = (1.0 / pres.slide_size.size.height) * desiredY


        # Create a full scale image
        bmp = sld.get_thumbnail(ScaleX, ScaleY)

        # save the image to disk in JPEG format
        print(out_dir)
        bmp.save(out_dir, draw.imaging.ImageFormat.png)
        return bmp
    
    
def create_pptx(presentation, page, result, doc, path2save, boxes = [], clases = [], num2class= {}): 

    res = result.pages[0]
    slide_layout = presentation.slide_layouts[6]
    presentation.slides.add_slide(slide_layout)
                                                
    raw_data = res.export()
    
    img = doc[0]
    
    height, width, _ = img.shape
    
    pixels_per_inch = 70  # Standard resolution for screens
    new_width = int(width * 914400 / pixels_per_inch)
    new_height = int(height * 914400/ pixels_per_inch)

    # Resize the slide layout
    presentation.slide_width = new_width
    presentation.slide_height = new_height
    
    for (box, cls) in zip(boxes, clases):
        slide = presentation.slides[page]
        x1, y1, x2, y2 = box
        
        x1 = x1 * new_width
        x2 = x2 * new_width
        
        y1 = y1 * new_height
        y2 = y2 * new_height
                    
        text = num2class[int(cls)]
        text = text.upper()
        txBox = slide.shapes.add_textbox(x1, y1, x2-x1, y2-y1)
        line = txBox.line
        line.width = Pt(5)
        if text == "STAMP":
            color = RGBColor(0, 255, 0)  # Red color
        if text == "SIGNATURE":
            color = RGBColor(255, 0, 0)
        if text == "QR":
            color = RGBColor(0, 0, 255)
        if text == "BARCODE":
            color = RGBColor(255, 0, 255)
        line.color.rgb = color 
        
        tf = txBox.text_frame

        tf.paragraphs[0].font.size = Pt(30)
        tf.paragraphs[0].text = text
        tf.paragraphs[0].font.color.rgb = color
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
    
    for block in raw_data["blocks"]:
        for line in block["lines"]:
            [x1, y1], [x2, y2] = line["geometry"]
            
            x1, x2 = x1*width, x2*width
            y1, y2 = y1*height, y2*height
            
            
            text = " ".join([word["value"] for word in line["words"]])

            add_line_in_bounding_boxes(presentation, page, text, [Pt(x1), Pt(y1), Pt(x2-x1), Pt(y2-y1)])

    presentation.save(path2save)
    

    
    
    
    
    
    
    
    
"""
RESTORMER UTILS 

"""
    
import numpy as np
import cv2
import torch
import torch.nn.functional as F
import torchvision.transforms.functional as TF
from runpy import run_path
from skimage import img_as_ubyte
import os

def save_img(filepath, img):
    os.makedirs(os.path.dirname(filepath), exist_ok=True)
    cv2.imwrite(filepath,cv2.cvtColor(img, cv2.COLOR_RGB2BGR))

def load_img(filepath):
    img = cv2.imread(filepath)
    img = cv2.cvtColor(img, cv2.COLOR_BGR2RGB)
    return img

def transform_output(output):
        # given a list of vector of size HxWx6 return a list of 2 tensors of size HxWx3
        if isinstance(output, list):
            pred_l0 = []
            pred_l1 = []
            for out in output: 
                pred_l0.append(out[:, :3, :])
                pred_l1.append(out[:, 3:, :])
            return pred_l0, pred_l1
        else: 
            return output[:, :3, :], output[:, 3:, :]

def separate_layers(file_, out_dir, device):
    # Load model architecture based on selected model
    parameters = {'inp_channels':3, 'out_channels':3, 'dim':48, 'num_blocks':[4,6,6,8], 'num_refinement_blocks':4, 'heads':[1,2,4,8], 'ffn_expansion_factor':2.66, 'bias':False, 'LayerNorm_type':'BiasFree', 'dual_pixel_task':False}
    load_arch = run_path("/hhome/ps2g07/document_analysis/github/Project_Synthesis2-/Code/Denoise/Restormer/basicsr/models/archs/restormer_6out_arch.py")
    model = load_arch['Restormer'](**parameters)
    model.to(device)

    weights = "/hhome/ps2g07/document_analysis/github/Project_Synthesis2-/Code/Denoise/Restormer/experiments/DocumentDenoisingLayers_ds_v4/models/zbest_psnr_l1.pth"
    checkpoint = torch.load(weights)
    model.load_state_dict(checkpoint['params'])
    model.eval()
    
    img_multiple_of = 8

    with torch.no_grad():
        if torch.cuda.is_available():
            torch.cuda.ipc_collect()
            torch.cuda.empty_cache()

        input_ = load_img(file_)
        
        h, w = input_.shape[:2]
        if h > w:
            input_ = np.array(TF.resize(TF.to_pil_image(input_), (1000, int(1000*w/h))))
        else:
            input_ = np.array(TF.resize(TF.to_pil_image(input_), (int(1000*h/w), 1000)))
        input_ = torch.from_numpy(input_).float().div(255.).permute(2,0,1).unsqueeze(0).to(device)
        
        # Pad the input if not_multiple_of 8
        height, width = input_.shape[2], input_.shape[3]
        H, W = ((height+img_multiple_of)//img_multiple_of)*img_multiple_of, ((width+img_multiple_of)//img_multiple_of)*img_multiple_of
        padh = H-height if height%img_multiple_of!=0 else 0
        padw = W-width if width%img_multiple_of!=0 else 0
        input_ = F.pad(input_, (0,padw,0,padh), 'reflect')

        restored = model(input_)
        restored_l0, restored_l1 = transform_output(restored)
        
        restored_l0 = torch.clamp(restored_l0, 0, 1)
        restored_l1 = torch.clamp(restored_l1, 0, 1)

        # Unpad the output
        restored_l0 = restored_l0[:,:,:height,:width]
        restored_l1 = restored_l1[:,:,:height,:width]

        restored_l0 = restored_l0.permute(0, 2, 3, 1).cpu().detach().numpy()
        restored_l1 = restored_l1.permute(0, 2, 3, 1).cpu().detach().numpy()
        
        restored_l0 = img_as_ubyte(restored_l0[0])
        restored_l1 = img_as_ubyte(restored_l1[0])

        f = os.path.splitext(os.path.split(file_)[-1])[0]
        # save_img((os.path.join(out_dir, f+'_l0.png')), restored_l0)
        save_img((os.path.join(out_dir, f+'_l0.png')), restored_l0)
        save_img((os.path.join(out_dir, f+'_l1.png')), restored_l1)

        return restored_l0, restored_l1
    
from pdf2image import convert_from_path
def pdf2png(pdf_path, png_root): 
    pages = convert_from_path(pdf_path, 500)
    os.makedirs(png_root, exist_ok=True)
    
    pdf_name = os.path.basename(pdf_path)
    for i, page in enumerate(pages):
        page.save(os.path.join(png_root, pdf_name.split(".")[0] + "_" + str(i) + ".png"), "PNG")
    
    return len(pages)


import img2pdf
from PIL import Image
def png2pdf(list_of_png_paths, pdf_root):    
    # converting into chunks using img2pdf
    a4inpt = (img2pdf.mm_to_pt(210),img2pdf.mm_to_pt(297))
    layout_fun = img2pdf.get_layout_fun(a4inpt)

    pdf_bytes = img2pdf.convert(list_of_png_paths, layout_fun=layout_fun)
    
    # writing to pdf file
    with open(pdf_root + "/layer_sep" + ".pdf", "wb") as f:
        f.write(pdf_bytes)
    
    
    
    
''' YOLO UTILS'''

import sys
sys.path.append('/hhome/ps2g07/document_analysis/github/Project_Synthesis2-/Nil/yolo_dir/ultralytics')
from ultralytics import YOLO
from pathlib import Path
import torch
import os

pretrained_weights = Path('/hhome/ps2g07/runs/detect/train31/weights/best.pt')

def infer_yolo(imgs, out_name):
    imgs = [imgs]
    width, height = Image.open(imgs[0]).size
    model = YOLO(pretrained_weights)  # load pre trained model
    results = model(imgs[0])
    idx_treshold = torch.where(results[0].boxes.conf > 0.6)[0]
    cls = results[0].boxes[idx_treshold].cls.unsqueeze(1)
    conf = results[0].boxes[idx_treshold].conf.unsqueeze(1)
    box = results[0].boxes[idx_treshold].xyxy
    updated_boxes = torch.cat((box, conf, cls), dim=1)
    results[0].update(boxes=updated_boxes)
    results[0].save(filename=out_name)
    # get for each detected thing the class and the bounding box
    print(results[0].boxes.xyxy)
    print(results[0].boxes.cls)
    print(results[0].names)
    
    results[0].boxes.xyxy[:, [0, 2]] = results[0].boxes.xyxy[:, [0, 2]] / width
    results[0].boxes.xyxy[:, [1, 3]] = results[0].boxes.xyxy[:, [1, 3]] / height
    
    print("\n"*5)
    
    return results[0].boxes.xyxy, results[0].boxes.cls, results[0].names
    
    
''' TSR UTILS'''

from collections import OrderedDict, defaultdict
import json
import argparse
import sys
import xml.etree.ElementTree as ET
import os
import random
import io

import torch
from torchvision import transforms
from PIL import Image
from fitz import Rect
import numpy as np
import pandas as pd
import matplotlib
#matplotlib.use('TkAgg')
import matplotlib.pyplot as plt
import matplotlib.patches as patches
from matplotlib.patches import Patch

# from main import get_model
sys.path.append("/hhome/ps2g07/document_analysis/github/Project_Synthesis2-/Code/TSR/table-transformer/src")
import postprocess
sys.path.append("/hhome/ps2g07/document_analysis/github/Project_Synthesis2-/Code/TSR/table-transformer/detr")
from models import build_model

Image.MAX_IMAGE_PIXELS = 200000000
                       #  192782800

class MaxResize(object):
   def __init__(self, max_size=800):
       self.max_size = max_size

   def __call__(self, image):
       width, height = image.size
       current_max_size = max(width, height)
       scale = self.max_size / current_max_size
       resized_image = image.resize((int(round(scale*width)), int(round(scale*height))))
       
       return resized_image

detection_transform = transforms.Compose([
   MaxResize(800),
   transforms.ToTensor(),
   transforms.Normalize([0.485, 0.456, 0.406], [0.229, 0.224, 0.225])
])

structure_transform = transforms.Compose([
   MaxResize(1000),
   transforms.ToTensor(),
   transforms.Normalize([0.485, 0.456, 0.406], [0.229, 0.224, 0.225])
])

def get_class_map(data_type):
   if data_type == 'structure':
       class_map = {
           'table': 0,
           'table column': 1,
           'table row': 2,
           'table column header': 3,
           'table projected row header': 4,
           'table spanning cell': 5,
           'no object': 6
       }
   elif data_type == 'detection':
       class_map = {'table': 0, 'table rotated': 1, 'no object': 2}
   return class_map

detection_class_thresholds = {
   "table": 0.5,
   "table rotated": 0.5,
   "no object": 10
}

structure_class_thresholds = {
   "table": 0.5,
   "table column": 0.5,
   "table row": 0.5,
   "table column header": 0.5,
   "table projected row header": 0.5,
   "table spanning cell": 0.5,
   "no object": 10
}


# for output bounding box post-processing
def box_cxcywh_to_xyxy(x):
   x_c, y_c, w, h = x.unbind(-1)
   b = [(x_c - 0.5 * w), (y_c - 0.5 * h), (x_c + 0.5 * w), (y_c + 0.5 * h)]
   return torch.stack(b, dim=1)

def rescale_bboxes(out_bbox, size):
   img_w, img_h = size
   b = box_cxcywh_to_xyxy(out_bbox)
   b = b * torch.tensor([img_w, img_h, img_w, img_h], dtype=torch.float32)
   return b

def iob(bbox1, bbox2):
   """
   Compute the intersection area over box area, for bbox1.
   """
   intersection = Rect(bbox1).intersect(bbox2)
   
   bbox1_area = Rect(bbox1).get_area()
   if bbox1_area > 0:
       return intersection.get_area() / bbox1_area
   
   return 0

def align_headers(headers, rows):
   """
   Adjust the header boundary to be the convex hull of the rows it intersects
   at least 50% of the height of.

   For now, we are not supporting tables with multiple headers, so we need to
   eliminate anything besides the top-most header.
   """
   
   aligned_headers = []

   for row in rows:
       row['column header'] = False

   header_row_nums = []
   for header in headers:
       for row_num, row in enumerate(rows):
           row_height = row['bbox'][3] - row['bbox'][1]
           min_row_overlap = max(row['bbox'][1], header['bbox'][1])
           max_row_overlap = min(row['bbox'][3], header['bbox'][3])
           overlap_height = max_row_overlap - min_row_overlap
           if overlap_height / row_height >= 0.5:
               header_row_nums.append(row_num)

   if len(header_row_nums) == 0:
       return aligned_headers

   header_rect = Rect()
   if header_row_nums[0] > 0:
       header_row_nums = list(range(header_row_nums[0]+1)) + header_row_nums

   last_row_num = -1
   for row_num in header_row_nums:
       if row_num == last_row_num + 1:
           row = rows[row_num]
           row['column header'] = True
           header_rect = header_rect.include_rect(row['bbox'])
           last_row_num = row_num
       else:
           # Break as soon as a non-header row is encountered.
           # This ignores any subsequent rows in the table labeled as a header.
           # Having more than 1 header is not supported currently.
           break

   header = {'bbox': list(header_rect)}
   aligned_headers.append(header)

   return aligned_headers

def refine_table_structure(table_structure, class_thresholds):
   """
   Apply operations to the detected table structure objects such as
   thresholding, NMS, and alignment.
   """
   rows = table_structure["rows"]
   columns = table_structure['columns']

   # Process the headers
   column_headers = table_structure['column headers']
   column_headers = postprocess.apply_threshold(column_headers, class_thresholds["table column header"])
   column_headers = postprocess.nms(column_headers)
   column_headers = align_headers(column_headers, rows)

   # Process spanning cells
   spanning_cells = [elem for elem in table_structure['spanning cells'] if not elem['projected row header']]
   projected_row_headers = [elem for elem in table_structure['spanning cells'] if elem['projected row header']]
   spanning_cells = postprocess.apply_threshold(spanning_cells, class_thresholds["table spanning cell"])
   projected_row_headers = postprocess.apply_threshold(projected_row_headers,
                                                       class_thresholds["table projected row header"])
   spanning_cells += projected_row_headers
   # Align before NMS for spanning cells because alignment brings them into agreement
   # with rows and columns first; if spanning cells still overlap after this operation,
   # the threshold for NMS can basically be lowered to just above 0
   spanning_cells = postprocess.align_supercells(spanning_cells, rows, columns)
   spanning_cells = postprocess.nms_supercells(spanning_cells)

   postprocess.header_supercell_tree(spanning_cells)

   table_structure['columns'] = columns
   table_structure['rows'] = rows
   table_structure['spanning cells'] = spanning_cells
   table_structure['column headers'] = column_headers

   return table_structure

def outputs_to_objects(outputs, img_size, class_idx2name):
    # print("class_idx2name", class_idx2name)
    m = outputs['pred_logits'].softmax(-1).max(-1)
    pred_labels = list(m.indices.detach().cpu().numpy())[0]
    pred_scores = list(m.values.detach().cpu().numpy())[0]
    pred_bboxes = outputs['pred_boxes'].detach().cpu()[0]
    pred_bboxes = [elem.tolist() for elem in rescale_bboxes(pred_bboxes, img_size)]

    # print("labels", pred_labels)
    # print("possible_labels", class_idx2name)
    objects = []
    for label, score, bbox in zip(pred_labels, pred_scores, pred_bboxes):
        class_label = class_idx2name[int(label)]
        # if class_label == "table": 
        if not class_label == 'no object':
            objects.append({'label': class_label, 'score': float(score),
                            'bbox': [float(elem) for elem in bbox]})
    # print("objects", objects)
    return objects

def objects_to_crops(img, tokens, objects, class_thresholds, padding=10):
   """
   Process the bounding boxes produced by the table detection model into
   cropped table images and cropped tokens.
   """

   table_crops = []
   for obj in objects:
       if obj['score'] < class_thresholds[obj['label']]:
           continue

       cropped_table = {}

       bbox = obj['bbox']
       bbox = [bbox[0]-padding, bbox[1]-padding, bbox[2]+padding, bbox[3]+padding]

       cropped_img = img.crop(bbox)

       table_tokens = [token for token in tokens if iob(token['bbox'], bbox) >= 0.5]
       for token in table_tokens:
           token['bbox'] = [token['bbox'][0]-bbox[0],
                            token['bbox'][1]-bbox[1],
                            token['bbox'][2]-bbox[0],
                            token['bbox'][3]-bbox[1]]

       # If table is predicted to be rotated, rotate cropped image and tokens/words:
       if obj['label'] == 'table rotated':
           cropped_img = cropped_img.rotate(270, expand=True)
           for token in table_tokens:
               bbox = token['bbox']
               bbox = [cropped_img.size[0]-bbox[3]-1,
                       bbox[0],
                       cropped_img.size[0]-bbox[1]-1,
                       bbox[2]]
               token['bbox'] = bbox
               

       cropped_table['image'] = cropped_img
       cropped_table['tokens'] = table_tokens
       cropped_table['bbox'] = bbox
       cropped_table['label'] = obj['label']

       table_crops.append(cropped_table)
#    print("table_crops", table_crops)
   return table_crops

def objects_to_structures(objects, tokens, class_thresholds):
   """
   Process the bounding boxes produced by the table structure recognition model into
   a *consistent* set of table structures (rows, columns, spanning cells, headers).
   This entails resolving conflicts/overlaps, and ensuring the boxes meet certain alignment
   conditions (for example: rows should all have the same width, etc.).
   """

   tables = [obj for obj in objects if obj['label'] == 'table']
   table_structures = []

   for table in tables:
       table_objects = [obj for obj in objects if iob(obj['bbox'], table['bbox']) >= 0.5]
       table_tokens = [token for token in tokens if iob(token['bbox'], table['bbox']) >= 0.5]
       
       structure = {}

       columns = [obj for obj in table_objects if obj['label'] == 'table column']
       rows = [obj for obj in table_objects if obj['label'] == 'table row']
       column_headers = [obj for obj in table_objects if obj['label'] == 'table column header']
       spanning_cells = [obj for obj in table_objects if obj['label'] == 'table spanning cell']
       for obj in spanning_cells:
           obj['projected row header'] = False
       projected_row_headers = [obj for obj in table_objects if obj['label'] == 'table projected row header']
       for obj in projected_row_headers:
           obj['projected row header'] = True
       spanning_cells += projected_row_headers
       for obj in rows:
           obj['column header'] = False
           for header_obj in column_headers:
               if iob(obj['bbox'], header_obj['bbox']) >= 0.5:
                   obj['column header'] = True

       # Refine table structures
       rows = postprocess.refine_rows(rows, table_tokens, class_thresholds['table row'])
       columns = postprocess.refine_columns(columns, table_tokens, class_thresholds['table column'])

       # Shrink table bbox to just the total height of the rows
       # and the total width of the columns
       row_rect = Rect()
       for obj in rows:
           row_rect.include_rect(obj['bbox'])
       column_rect = Rect() 
       for obj in columns:
           column_rect.include_rect(obj['bbox'])
       table['row_column_bbox'] = [column_rect[0], row_rect[1], column_rect[2], row_rect[3]]
       table['bbox'] = table['row_column_bbox']

       # Process the rows and columns into a complete segmented table
       columns = postprocess.align_columns(columns, table['row_column_bbox'])
       rows = postprocess.align_rows(rows, table['row_column_bbox'])

       structure['rows'] = rows
       structure['columns'] = columns
       structure['column headers'] = column_headers
       structure['spanning cells'] = spanning_cells

       if len(rows) > 0 and len(columns) > 1:
           structure = refine_table_structure(structure, class_thresholds)

       table_structures.append(structure)

   return table_structures

def structure_to_cells(table_structure, tokens):
   """
   Assuming the row, column, spanning cell, and header bounding boxes have
   been refined into a set of consistent table structures, process these
   table structures into table cells. This is a universal representation
   format for the table, which can later be exported to Pandas or CSV formats.
   Classify the cells as header/access cells or data cells
   based on if they intersect with the header bounding box.
   """
   columns = table_structure['columns']
   rows = table_structure['rows']
   spanning_cells = table_structure['spanning cells']
   cells = []
   subcells = []

   # Identify complete cells and subcells
   for column_num, column in enumerate(columns):
       for row_num, row in enumerate(rows):
           column_rect = Rect(list(column['bbox']))
           row_rect = Rect(list(row['bbox']))
           cell_rect = row_rect.intersect(column_rect)
           header = 'column header' in row and row['column header']
           cell = {'bbox': list(cell_rect), 'column_nums': [column_num], 'row_nums': [row_num],
                   'column header': header}

           cell['subcell'] = False
           for spanning_cell in spanning_cells:
               spanning_cell_rect = Rect(list(spanning_cell['bbox']))
               if (spanning_cell_rect.intersect(cell_rect).get_area()
                       / cell_rect.get_area()) > 0.5:
                   cell['subcell'] = True
                   break

           if cell['subcell']:
               subcells.append(cell)
           else:
               #cell text = extract_text_inside_bbox(table_spans, cell['bbox'])
               #cell['cell text'] = cell text
               cell['projected row header'] = False
               cells.append(cell)

   for spanning_cell in spanning_cells:
       spanning_cell_rect = Rect(list(spanning_cell['bbox']))
       cell_columns = set()
       cell_rows = set()
       cell_rect = None
       header = True
       for subcell in subcells:
           subcell_rect = Rect(list(subcell['bbox']))
           subcell_rect_area = subcell_rect.get_area()
           if (subcell_rect.intersect(spanning_cell_rect).get_area()
                   / subcell_rect_area) > 0.5:
               if cell_rect is None:
                   cell_rect = Rect(list(subcell['bbox']))
               else:
                   cell_rect.include_rect(Rect(list(subcell['bbox'])))
               cell_rows = cell_rows.union(set(subcell['row_nums']))
               cell_columns = cell_columns.union(set(subcell['column_nums']))
               # By convention here, all subcells must be classified
               # as header cells for a spanning cell to be classified as a header cell;
               # otherwise, this could lead to a non-rectangular header region
               header = header and 'column header' in subcell and subcell['column header']
       if len(cell_rows) > 0 and len(cell_columns) > 0:
           cell = {'bbox': list(cell_rect), 'column_nums': list(cell_columns), 'row_nums': list(cell_rows),
                   'column header': header, 'projected row header': spanning_cell['projected row header']}
           cells.append(cell)

   # Compute a confidence score based on how well the page tokens
   # slot into the cells reported by the model
   _, _, cell_match_scores = postprocess.slot_into_containers(cells, tokens)
   try:
       mean_match_score = sum(cell_match_scores) / len(cell_match_scores)
       min_match_score = min(cell_match_scores)
       confidence_score = (mean_match_score + min_match_score)/2
   except:
       confidence_score = 0

   # Dilate rows and columns before final extraction
   #dilated_columns = fill_column_gaps(columns, table_bbox)
   dilated_columns = columns
   #dilated_rows = fill_row_gaps(rows, table_bbox)
   dilated_rows = rows
   for cell in cells:
       column_rect = Rect()
       for column_num in cell['column_nums']:
           column_rect.include_rect(list(dilated_columns[column_num]['bbox']))
       row_rect = Rect()
       for row_num in cell['row_nums']:
           row_rect.include_rect(list(dilated_rows[row_num]['bbox']))
       cell_rect = column_rect.intersect(row_rect)
       cell['bbox'] = list(cell_rect)

   span_nums_by_cell, _, _ = postprocess.slot_into_containers(cells, tokens, overlap_threshold=0.001,
                                                              unique_assignment=True, forced_assignment=False)

   for cell, cell_span_nums in zip(cells, span_nums_by_cell):
       cell_spans = [tokens[num] for num in cell_span_nums]
       # TODO: Refine how text is extracted; should be character-based, not span-based;
       # but need to associate 
       cell['cell text'] = postprocess.extract_text_from_spans(cell_spans, remove_integer_superscripts=False)
       cell['spans'] = cell_spans
       
   # Adjust the row, column, and cell bounding boxes to reflect the extracted text
   num_rows = len(rows)
   rows = postprocess.sort_objects_top_to_bottom(rows)
   num_columns = len(columns)
   columns = postprocess.sort_objects_left_to_right(columns)
   min_y_values_by_row = defaultdict(list)
   max_y_values_by_row = defaultdict(list)
   min_x_values_by_column = defaultdict(list)
   max_x_values_by_column = defaultdict(list)
   for cell in cells:
       min_row = min(cell["row_nums"])
       max_row = max(cell["row_nums"])
       min_column = min(cell["column_nums"])
       max_column = max(cell["column_nums"])
       for span in cell['spans']:
           min_x_values_by_column[min_column].append(span['bbox'][0])
           min_y_values_by_row[min_row].append(span['bbox'][1])
           max_x_values_by_column[max_column].append(span['bbox'][2])
           max_y_values_by_row[max_row].append(span['bbox'][3])
   for row_num, row in enumerate(rows):
       if len(min_x_values_by_column[0]) > 0:
           row['bbox'][0] = min(min_x_values_by_column[0])
       if len(min_y_values_by_row[row_num]) > 0:
           row['bbox'][1] = min(min_y_values_by_row[row_num])
       if len(max_x_values_by_column[num_columns-1]) > 0:
           row['bbox'][2] = max(max_x_values_by_column[num_columns-1])
       if len(max_y_values_by_row[row_num]) > 0:
           row['bbox'][3] = max(max_y_values_by_row[row_num])
   for column_num, column in enumerate(columns):
       if len(min_x_values_by_column[column_num]) > 0:
           column['bbox'][0] = min(min_x_values_by_column[column_num])
       if len(min_y_values_by_row[0]) > 0:
           column['bbox'][1] = min(min_y_values_by_row[0])
       if len(max_x_values_by_column[column_num]) > 0:
           column['bbox'][2] = max(max_x_values_by_column[column_num])
       if len(max_y_values_by_row[num_rows-1]) > 0:
           column['bbox'][3] = max(max_y_values_by_row[num_rows-1])
   for cell in cells:
       row_rect = Rect()
       column_rect = Rect()
       for row_num in cell['row_nums']:
           row_rect.include_rect(list(rows[row_num]['bbox']))
       for column_num in cell['column_nums']:
           column_rect.include_rect(list(columns[column_num]['bbox']))
       cell_rect = row_rect.intersect(column_rect)
       if cell_rect.get_area() > 0:
           cell['bbox'] = list(cell_rect)
           pass

   return cells, confidence_score

def visualize_detected_tables(img, det_tables, out_path):
   
   plt.imshow(img, interpolation="lanczos")
   plt.gcf().set_size_inches(20, 20)
   ax = plt.gca()
   
#    print("det_tables", det_tables)
   for det_table in det_tables:
       bbox = det_table['bbox']
       if det_table['label'] == 'table':
           facecolor = (1, 0, 0.45)
           edgecolor = (1, 0, 0.45)
           alpha = 0.3
           linewidth = 2
           hatch='//////'
       elif det_table['label'] == 'table rotated':
           facecolor = (0.95, 0.6, 0.1)
           edgecolor = (0.95, 0.6, 0.1)
           alpha = 0.3
           linewidth = 2
           hatch='//////'
       else:
           continue

       rect = patches.Rectangle(bbox[:2], bbox[2]-bbox[0], bbox[3]-bbox[1], linewidth=linewidth, 
                                   edgecolor='none',facecolor=facecolor, alpha=0.1)
       ax.add_patch(rect)
       rect = patches.Rectangle(bbox[:2], bbox[2]-bbox[0], bbox[3]-bbox[1], linewidth=linewidth, 
                                   edgecolor=edgecolor,facecolor='none',linestyle='-', alpha=alpha)
       ax.add_patch(rect)
       rect = patches.Rectangle(bbox[:2], bbox[2]-bbox[0], bbox[3]-bbox[1], linewidth=0, 
                                   edgecolor=edgecolor,facecolor='none',linestyle='-', hatch=hatch, alpha=0.2)
       ax.add_patch(rect)

   plt.xticks([], [])
   plt.yticks([], [])

   legend_elements = [Patch(facecolor=(1, 0, 0.45), edgecolor=(1, 0, 0.45),
                               label='Table', hatch='//////', alpha=0.3),
                       Patch(facecolor=(0.95, 0.6, 0.1), edgecolor=(0.95, 0.6, 0.1),
                               label='Table (rotated)', hatch='//////', alpha=0.3)]
   plt.legend(handles=legend_elements, bbox_to_anchor=(0.5, -0.02), loc='upper center', borderaxespad=0,
                   fontsize=10, ncol=2)  
   plt.gcf().set_size_inches(10, 10)
   plt.axis('off')
   plt.savefig(out_path, bbox_inches='tight', dpi=150)
   plt.close()

   return

def visualize_cells(img, allstructure, out_path, alltables):
#    print(table.keys())
    plt.imshow(img, interpolation="lanczos")
    plt.gcf().set_size_inches(20, 20)
    ax = plt.gca()
    
    for cells, table in zip(allstructure, alltables): 
        for cell in cells: 
            bbox = cell['bbox']
            table_bbox = table['bbox']
            padding= 20
            
            bbox = [bbox[0]+table_bbox[0]-padding, bbox[1]+table_bbox[1]-padding,
                    bbox[2]+table_bbox[0]+padding, bbox[3]+table_bbox[1]+padding]

            if cell['column header']:
                facecolor = (1, 0, 0.45)
                edgecolor = (1, 0, 0.45)
                alpha = 0.3
                linewidth = 2
                hatch='//////'
            elif cell['projected row header']:
                facecolor = (0.95, 0.6, 0.1)
                edgecolor = (0.95, 0.6, 0.1)
                alpha = 0.3
                linewidth = 2
                hatch='//////'
            else:
                facecolor = (0.3, 0.74, 0.8)
                edgecolor = (0.3, 0.7, 0.6)
                alpha = 0.3
                linewidth = 2
                hatch='\\\\\\\\\\\\'

            rect = patches.Rectangle(bbox[:2], bbox[2]-bbox[0], bbox[3]-bbox[1], linewidth=linewidth, 
                                        edgecolor='none',facecolor=facecolor, alpha=0.1)
            ax.add_patch(rect)
            rect = patches.Rectangle(bbox[:2], bbox[2]-bbox[0], bbox[3]-bbox[1], linewidth=linewidth, 
                                        edgecolor=edgecolor,facecolor='none',linestyle='-', alpha=alpha)
            ax.add_patch(rect)
            rect = patches.Rectangle(bbox[:2], bbox[2]-bbox[0], bbox[3]-bbox[1], linewidth=0, 
                                        edgecolor=edgecolor,facecolor='none',linestyle='-', hatch=hatch, alpha=0.2)
            ax.add_patch(rect)

    plt.xticks([], [])
    plt.yticks([], [])

    legend_elements = [Patch(facecolor=(0.3, 0.74, 0.8), edgecolor=(0.3, 0.7, 0.6),
                                label='Data cell', hatch='\\\\\\\\\\\\', alpha=0.3),
                        Patch(facecolor=(1, 0, 0.45), edgecolor=(1, 0, 0.45),
                                label='Column header cell', hatch='//////', alpha=0.3),
                        Patch(facecolor=(0.95, 0.6, 0.1), edgecolor=(0.95, 0.6, 0.1),
                                label='Projected row header cell', hatch='//////', alpha=0.3)]
    plt.legend(handles=legend_elements, bbox_to_anchor=(0.5, -0.02), loc='upper center', borderaxespad=0,
                fontsize=10, ncol=3)  
   
    plt.gcf().set_size_inches(10, 10)
    plt.axis('off')
    plt.savefig(out_path, bbox_inches='tight', dpi=150)
    plt.close()

class TableExtractionPipeline(object):
   def __init__(self, device=None,
                det_model_path=None, str_model_path=None,
                det_config_path=None, str_config_path=None):

       self.det_device = device
       self.str_device = device

       self.det_class_name2idx = get_class_map('detection')
       self.det_class_idx2name = {v:k for k, v in self.det_class_name2idx.items()}
       self.det_class_thresholds = detection_class_thresholds

       self.str_class_name2idx = get_class_map('structure')
       self.str_class_idx2name = {v:k for k, v in self.str_class_name2idx.items()}
       self.str_class_thresholds = structure_class_thresholds

       if not det_config_path is None:
           with open(det_config_path, 'r') as f:
               det_config = json.load(f)
           det_args = type('Args', (object,), det_config)
           det_device = device
           self.det_model, _, _ = build_model(det_args)
           # print("Detection model initialized.")

           if not det_model_path is None:
               self.det_model.load_state_dict(torch.load(det_model_path,
                                                    map_location=torch.device(device)))
               self.det_model.to(device)
               self.det_model.eval()
               # print("Detection model weights loaded.")
           else:
               self.det_model = None

       if not str_config_path is None:
           with open(str_config_path, 'r') as f:
               str_config = json.load(f)
           str_args = type('Args', (object,), str_config)
           str_device = device
           self.str_model, _, _ = build_model(str_args)
           # print("Structure model initialized.")

           if not str_model_path is None:
               self.str_model.load_state_dict(torch.load(str_model_path,
                                                    map_location=torch.device(device)))
               self.str_model.to(device)
               self.str_model.eval()
               # print("Structure model weights loaded.")ç
           else:
               self.str_model = None


   def __call__(self, page_image, page_tokens=None):
       return self.extract(self, page_image, page_tokens)

   def detect(self, img, tokens=None, out_objects=True, out_crops=False, crop_padding=10):
       out_formats = {}
       if self.det_model is None:
           print("No detection model loaded.")
           return out_formats

       # Transform the image how the model expects it
       img_tensor = detection_transform(img)

       # Run input image through the model
       outputs = self.det_model([img_tensor.to(self.det_device)])
    #    print("classes", self.det_class_idx2name)

       # Post-process detected objects, assign class labels
       objects = outputs_to_objects(outputs, img.size, self.det_class_idx2name)
       if out_objects:
           out_formats['objects'] = objects
       if not out_crops:
           return out_formats

       # print(objects[0]['bbox'])
       # Crop image and tokens for detected table
       if out_crops:
            tables_crops = objects_to_crops(img, tokens, objects, self.det_class_thresholds,
                                           padding=crop_padding)
            out_formats['crops'] = tables_crops
            # print("out_formats", out_formats)
       return out_formats

   def recognize(self, img, tokens=None, out_objects=False, out_cells=False,
                 out_html=False, out_csv=False):
       out_formats = {}
       if self.str_model is None:
           print("No structure model loaded.")
           return out_formats

       if not (out_objects or out_cells or out_html or out_csv):
           print("No output format specified")
           return out_formats

       # Transform the image how the model expects it
       img_tensor = structure_transform(img)

       # Run input image through the model
       outputs = self.str_model([img_tensor.to(self.str_device)])
        
       # Post-process detected objects, assign class labels
       objects = outputs_to_objects(outputs, img.size, self.str_class_idx2name)
    #    print("objects", objects)
    #    print("class_idx2name", self.str_class_idx2name)
       if out_objects:
           out_formats['objects'] = objects
       if not (out_cells or out_html or out_csv):
           return out_formats

       # Further process the detected objects so they correspond to a consistent table 
       tables_structure = objects_to_structures(objects, tokens, self.str_class_thresholds)

       # Enumerate all table cells: grid cells and spanning cells
       tables_cells = [structure_to_cells(structure, tokens)[0] for structure in tables_structure]
       if out_cells:
           out_formats['cells'] = tables_cells

       return out_formats

   def extract(self, img, tokens=None, out_objects=True, out_cells=False,
               out_html=False, out_csv=False, crop_padding=20):

       if tokens is None:
           tokens = []

       detect_out = self.detect(img, tokens = tokens, out_objects=False, out_crops=True, crop_padding=crop_padding)
       cropped_tables = detect_out['crops']

       # save in a json file the bbox of the detected tables img_file.replace(".jpg", "_{}.{}".format(idx, key))
       # with open('detected_tables.json', 'w') as f:
       #     json.dump(detect_out['bbox'][0], f)

       for table_idx, extracted_table in enumerate(cropped_tables):
            for key, val in extracted_table.items():
                # img = extracted_table['full_image']
                img_file = '{}.jpg'.format(table_idx)
                
                out_file = img_file.replace(".jpg", "_objects.json")
                with open(os.path.join(out_dir, out_file), 'w') as f:
                    json.dump(val, f)

                out_file = img_file.replace(".jpg", "_fig_tables.jpg")  # <- 
                out_path = os.path.join(out_dir, out_file)
                visualize_detected_tables(img, val, out_path)

       extracted_tables = []
       for table in cropped_tables:
           img = table['image']
           tokens = table['tokens']
           bbox = table['bbox']

           extracted_table = self.recognize(img, tokens=tokens, out_objects=out_objects,
                                      out_cells=out_cells, out_html=out_html, out_csv=out_csv)
           
           extracted_table['image'] = img
           extracted_table['tokens'] = tokens
           extracted_table['bbox_table'] = bbox
           extracted_tables.append(extracted_table)
        
    #    print("extracted_tables", extracted_tables)

       return extracted_tables

       
def infer_table(img_path, out_dir, device):
    detection_config_path = "/hhome/ps2g07/document_analysis/github/Project_Synthesis2-/Code/TSR/table-transformer/src/detection_config.json"
    detection_model_path = "/hhome/ps2g07/document_analysis/github/Project_Synthesis2-/Code/TSR/table-transformer/ckpts/pubtables1m_detection_detr_r18.pth"
    structure_config_path = "/hhome/ps2g07/document_analysis/github/Project_Synthesis2-/Code/TSR/table-transformer/src/structure_config.json"
    structure_model_path = "/hhome/ps2g07/document_analysis/github/Project_Synthesis2-/Code/TSR/table-transformer/ckpts/pubtables1m_structure_detr_r18.pth"
    
    img_file, img_root = img_path.split('/')[-1], '/'.join(img_path.split('/')[:-1])
    
    os.makedirs(out_dir, exist_ok=True)

    pipe = TableExtractionPipeline(device=device, 
                                    det_config_path=detection_config_path, det_model_path=detection_model_path, 
                                    str_config_path=structure_config_path, str_model_path=structure_model_path)

    # Load images
    img = Image.open(os.path.join(img_root, img_file))

    tokens = []
    detect_out = pipe.detect(img, tokens = tokens, out_objects=False, out_crops=True, crop_padding=10)
    cropped_tables = detect_out['crops']

    # print("cropped_tables", cropped_tables)
    # img_file = 'tables.jpg'
    
    out_file = "tables.json"
    cropped_wo_img = [{k: v for k, v in elem.items() if not k == 'image'} for elem in cropped_tables]
    with open(os.path.join(out_dir, out_file), 'w') as f:
        json.dump(cropped_wo_img, f)

    out_file = "tables.jpg"  # <- 
    out_path = os.path.join(out_dir, out_file)
    visualize_detected_tables(img, cropped_tables, out_path)

    # extracted_tables = []
    list_cells = []
    for table_idx, table in enumerate(cropped_tables):
        img_cropped = table['image']
        tokens = table['tokens']
        # bbox = table['bbox']

        extracted_table = pipe.recognize(img_cropped, tokens=tokens, out_objects=False,
                                    out_cells=True, out_html=False, out_csv=False)
        
        extracted_table['image'] = img_cropped
        extracted_table['tokens'] = tokens
        # extracted_tables.append(extracted_table)

        out_file = '{}_cells.json'.format(table_idx)
        
        with open(os.path.join(out_dir, out_file), 'w') as f:
            json.dump(extracted_table["cells"][0], f)
            
        
        list_cells.append(extracted_table["cells"][0])
        
    out_path = os.path.join(out_dir, f"cells.jpg")
    visualize_cells(img, list_cells, out_path, cropped_tables)


if __name__ == "__main__":
    detection_config_path = "detection_config.json"
    detection_model_path = "/hhome/ps2g07/document_analysis/github/Project_Synthesis2-/Code/TSR/table-transformer/ckpts/pubtables1m_detection_detr_r18.pth"
    structure_config_path = "structure_config.json"
    structure_model_path = "/hhome/ps2g07/document_analysis/github/Project_Synthesis2-/Code/TSR/table-transformer/ckpts/pubtables1m_structure_detr_r18.pth"
    
    device = "cuda"
    img_file = "38.jpg"
    out_dir = "/hhome/ps2g07/document_analysis/github/Project_Synthesis2-/Code/TSR/test_output"
    img_path = "/hhome/ps2g07/document_analysis/github/Project_Synthesis2-/Sample documents - JPG - NoName/38.jpg"

    infer_table(img_path, out_dir, device)