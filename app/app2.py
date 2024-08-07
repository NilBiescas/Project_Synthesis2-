from flask import Flask, render_template, request, send_file, send_from_directory, flash, redirect, url_for
from werkzeug.utils import secure_filename
import os
from helpers import *
from pptx.util import Pt
from pptx import Presentation
import torch
import aspose.slides as slides
import aspose.pydrawing as draw

from doctr import io
from doctr import utils
from time import sleep
import shutil
#from doctr.io import DocumentFile
#from doctr.utils.visualization import vis_and_synth


from PIL import Image
import matplotlib.pyplot as plt



app = Flask(__name__)

UPLOAD_FOLDER = 'uploads'
ALLOWED_EXTENSIONS = {'pdf', 'png', 'jpg', 'jpeg'}

device = "cuda" if torch.cuda.is_available() else "cpu"

app.config['UPLOAD_FOLDER'] = UPLOAD_FOLDER

def allowed_file(filename):
    return '.' in filename and filename.rsplit('.', 1)[1].lower() in ALLOWED_EXTENSIONS

def inference_models(file, device="cpu"): 
    file_name = secure_filename(file.filename)
    path2file = os.path.join(app.config['UPLOAD_FOLDER'], file_name)
    
    num_pages = pdf2png(path2file, app.config['UPLOAD_FOLDER'] + "/" + file_name[:-4])
    
    files = os.listdir(path2file[:-4])
    
    presentation = Presentation() 
        
    #infer_table(os.path.join(path2file[:-4], files[0]), out_dir = app.config['UPLOAD_FOLDER'], device=device)
    for page in range(num_pages):
        # infer_table(os.path.join(path2file[:-4], files[page]), out_dir = app.config['UPLOAD_FOLDER'], device=device)
        separate_layers(os.path.join(path2file[:-4], files[page]), out_dir = app.config['UPLOAD_FOLDER'], device=device)
        path2l0 = os.path.join(app.config['UPLOAD_FOLDER'], file_name[:-4]) + f"_{page}_l0.png"
        path2l1 = os.path.join(app.config['UPLOAD_FOLDER'], file_name[:-4] + f"_{page}_l1.png")
        result, doc = docTR_inference(path2l0, device=device)
    
        slides_pptx = file_name.replace(".pdf", f"{page}.pptx")
        path2save_slide = os.path.join(app.config['UPLOAD_FOLDER'], slides_pptx)
        file_name_png = file_name.split(".")[0] + f"_{page}_ocr.png"
        path2savePng = os.path.join(app.config['UPLOAD_FOLDER'], file_name_png)

        path2yolo = os.path.join(app.config['UPLOAD_FOLDER'], file_name[:-4] + f"{page}_yolo.png")
        file_name_yolo = file_name.split(".")[0] + f"{page}_yolo.png"
        boxes, clases, num2class = infer_yolo(path2l1, out_name=path2yolo)
        create_pptx(presentation, page, result, doc, path2save_slide, boxes, clases, num2class)
        
        plot_slide(path2save_slide, page, out_dir=path2savePng)
    return slides_pptx, file_name_png, file_name_yolo
    

@app.route('/') 
def index():
    return render_template('index.html')

@app.route('/upload', methods=['GET', 'POST'])
def upload_file():
    global name_file_uploaded 
    if request.method == 'POST':
        cleanup_upload_folder()
        if 'file' not in request.files:
            flash('No file part', 'error')
            return redirect(request.url)
        file = request.files['file']
        if file.filename == '':
            flash('No selected file', 'error')
            return render_template('upload.html')
        
        if file and allowed_file(file.filename):
            filename = secure_filename(file.filename)
            name_file_uploaded = filename
            file_path = os.path.join(app.config['UPLOAD_FOLDER'], filename)
            
            file.save(file_path)
            
            file_pptx, file_png, file_yolo = inference_models(file, device=device)
            
            # file_pptx = filename.replace(".pdf", ".pptx")
            # file_png = filename.replace(".pdf", "_0_ocr.png")
            # file_yolo = filename.replace(".pdf", "0_yolo.png")
            
            pptx_path = os.path.join(app.config['UPLOAD_FOLDER'], file_pptx)
            pdf_name = filename.replace(".pdf", "2.pdf")
            pdf_path = os.path.join(app.config['UPLOAD_FOLDER'], pdf_name)
            with slides.Presentation(pptx_path) as pres:
                pres.save(pdf_path, slides.export.SaveFormat.PDF)
            
            #sleep(10) 
            return render_template('preview.html', file_pptx=file_pptx, file_png=file_png, file_pdf = pdf_name, file_yolo=file_yolo)
        
    return render_template('upload.html')

@app.route('/contact')
def contact_us():
    return render_template('contact.html')

@app.route('/team')
def team():
    return render_template('team.html')

@app.route('/preview/<filename>', methods=['GET'])
def preview_file(filename):
    return send_from_directory(app.config['UPLOAD_FOLDER'], filename)

@app.route('/download/<filename>', methods=['GET'])
def download_file(filename):
    return send_file(os.path.join(app.config['UPLOAD_FOLDER'], filename), as_attachment=True)

@app.route('/uploads/<filename>', methods=['GET']) 
def uploaded_file(filename):
    if filename == "l0.png":
        name_file = name_file_uploaded.split(".")[0] + "_0_l0.png"
        return send_from_directory(app.config['UPLOAD_FOLDER'], name_file)
    elif filename == "l1.png":
        name_file = name_file_uploaded.split(".")[0] + "_0_l1.png"
        return send_from_directory(app.config['UPLOAD_FOLDER'], name_file)
    elif filename == "ocr.png":
        name_file = name_file_uploaded.split(".")[0] + "_0_ocr.png"
        return send_from_directory(app.config['UPLOAD_FOLDER'], name_file)
    elif filename == "yolo.png":
        name_file = name_file_uploaded.split(".")[0] + "_yolo.png"
        return send_from_directory(app.config['UPLOAD_FOLDER'], name_file)
    elif filename == "in.pdf":
        return send_from_directory(app.config['UPLOAD_FOLDER'], name_file_uploaded)
    return send_from_directory(app.config['UPLOAD_FOLDER'], filename)

def cleanup_upload_folder():
    # pass
    # remove every file and dir in the upload folder
    for root, dirs, files in os.walk(app.config['UPLOAD_FOLDER']):
        for file in files:
            os.remove(os.path.join(root, file))
        for dir in dirs:
            shutil.rmtree(os.path.join(root, dir))

import atexit
atexit.register(cleanup_upload_folder)

if __name__ == '__main__':
    app.run(debug=True, port=5002)